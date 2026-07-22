"""
Build PowerPoint from Deep-Read sections with extracted paper figures.
"""

from __future__ import annotations

from pathlib import Path
from typing import List, Optional

from etl.pdf_rich_loader import RichPDFDocument
from orchestration.deep_read.section_agents import SectionResult

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None


def _add_bullet_slide(prs, title: str, body: str, image_path: Optional[str] = None):
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title[:80]

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    # First 8 bullet lines from section
    lines = [ln.strip() for ln in body.split("\n") if ln.strip()][:8]
    for i, line in enumerate(lines):
        line = line.lstrip("-•* ").strip()
        if not line:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line[:300]
        p.level = 0
        p.font.size = Pt(14)

    if image_path and Path(image_path).exists():
        try:
            slide.shapes.add_picture(image_path, Inches(6.5), Inches(1.5), width=Inches(3.0))
        except Exception:
            pass


def build_ppt(
    sections: List[SectionResult],
    doc: RichPDFDocument,
    output_path: str,
) -> str:
    if Presentation is None:
        raise ImportError("python-pptx is required: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = doc.title[:100]
    subtitle = title_slide.placeholders[1]
    subtitle.text = f"Deep-Read Summary · {doc.total_pages} pages"

    # Map figures by page for method / figures slides
    figures_by_page = {}
    for fig in doc.figures:
        figures_by_page.setdefault(fig.page, []).append(fig.path)

    for sec in sections:
        img = None
        if sec.section_id == "method" and doc.figures:
            img = doc.figures[0].path
        elif sec.section_id == "figures_tables" and doc.figures:
            img = doc.figures[min(1, len(doc.figures) - 1)].path
        _add_bullet_slide(prs, sec.title, sec.content, image_path=img)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)
